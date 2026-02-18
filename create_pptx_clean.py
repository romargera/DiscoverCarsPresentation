import os
import glob
from pptx import Presentation
from pptx.util import Inches

# Configuration
ARTIFACTS_DIR = "/Users/roman/.gemini/antigravity/brain/6c4dd8ec-b63f-4c08-b470-63694ed95444"
OUTPUT_FILE = "/Users/roman/Documents/Pet Projects/DiscoverCarsPresentation/DiscoverCars_Presentation.pptx"

def create_presentation():
    # Create new presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9 (standard wide)
    # 13.333 inches x 7.5 inches = 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Find all slide screenshots
    # Pattern: slide_XX_*.png
    # Collect latest file for each slide pattern
    slide_files = []
    prefixes = [
        "slide_01_", "slide_02_", "slide_03_", "slide_04_", 
        "slide_05_", "slide_06_", "slide_07_", "slide_08_",
        "slide_09_", "slide_10_", "slide_11_", "slide_12_",
        "slide_13_", "slide_14_", "slide_15_"
    ]
    
    for prefix in prefixes:
        pattern = os.path.join(ARTIFACTS_DIR, prefix + "*.png")
        all_matches = glob.glob(pattern)
        # Filter strictly for this session's timestamp
        matches = sorted([f for f in all_matches if "_1771433" in f])
        
        if matches:
            # Pick the last one (latest timestamp)
            latest = matches[-1]
            slide_files.append(latest)
            print(f"Selected for {prefix}: {os.path.basename(latest)}")
        else:
            print(f"WARNING: No matching files found for {prefix}")
    
    print(f"Found {len(slide_files)} slides to include.")
    
    for image_path in slide_files:
        print(f"Adding slide from: {os.path.basename(image_path)}")
        
        # Add a blank slide (layout 6 is usually blank)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add the image to fill the slide
        # left, top, width, height
        slide.shapes.add_picture(
            image_path, 
            0, 0, 
            width=prs.slide_width, 
            height=prs.slide_height
        )

    # Save the presentation
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to: {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
